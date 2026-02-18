#!/usr/bin/env python3
"""Extract slide images from NotebookLM output (.pptx or .pdf).

Extracts embedded images from .pptx files (preferred) or renders
PDF pages as images. Output is organized for use in Reveal.js slides.

Usage:
    # From .pptx (extracts embedded images per slide):
    python3 extract_slides.py presentation.pptx -o output_dir/

    # From .pdf (renders each page as a full-slide image):
    python3 extract_slides.py slides.pdf -o output_dir/

    # Render full slides from .pptx (each slide as one image):
    python3 extract_slides.py presentation.pptx -o output_dir/ --full-slides
"""

import argparse
import os
import subprocess
import sys
from pathlib import Path

VENV_PYTHON = Path(__file__).parent / ".venv" / "bin" / "python3"


def extract_pptx_images(pptx_path, output_dir):
    """Extract embedded images from a .pptx file, organized by slide."""
    # Use venv python for python-pptx
    python = str(VENV_PYTHON) if VENV_PYTHON.exists() else sys.executable

    script = f'''
import sys
sys.path.insert(0, "")
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import json

prs = Presentation("{pptx_path}")
output_dir = Path("{output_dir}")
output_dir.mkdir(parents=True, exist_ok=True)

manifest = []

for slide_idx, slide in enumerate(prs.slides, 1):
    slide_title = ""
    slide_images = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and not slide_title:
                    slide_title = text
                    break

        if shape.shape_type == 13:  # Picture
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"

            filename = f"slide-{{slide_idx:02d}}-img-{{len(slide_images)+1}}.{{ext}}"
            filepath = output_dir / filename
            filepath.write_bytes(image.blob)
            slide_images.append(filename)

    manifest.append({{
        "slide": slide_idx,
        "title": slide_title,
        "images": slide_images,
    }})

# Write manifest
manifest_path = output_dir / "manifest.json"
manifest_path.write_text(json.dumps(manifest, indent=2))

total_images = sum(len(s["images"]) for s in manifest)
print(f"Extracted {{total_images}} images from {{len(manifest)}} slides")
print(f"Manifest: {{manifest_path}}")
'''

    result = subprocess.run([python, "-c", script], capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Error: {result.stderr}", file=sys.stderr)
        sys.exit(1)
    print(result.stdout, end="")


def render_pptx_full_slides(pptx_path, output_dir):
    """Render each .pptx slide as a full image using LibreOffice + pdftoppm."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Convert pptx to pdf first
    pdf_path = output_dir / "temp_slides.pdf"
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(output_dir), str(pptx_path)],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        print("LibreOffice not available for full-slide rendering.", file=sys.stderr)
        print("Install with: sudo apt install libreoffice", file=sys.stderr)
        sys.exit(1)

    # Find the converted PDF
    pdf_name = Path(pptx_path).stem + ".pdf"
    pdf_path = output_dir / pdf_name

    if pdf_path.exists():
        render_pdf_pages(str(pdf_path), str(output_dir))
        pdf_path.unlink()  # Clean up temp PDF
    else:
        print(f"PDF conversion failed: {pdf_path} not found", file=sys.stderr)


def render_pdf_pages(pdf_path, output_dir):
    """Render each PDF page as a high-res PNG image."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Try pdftoppm (poppler-utils)
    result = subprocess.run(
        ["pdftoppm", "-png", "-r", "300", pdf_path,
         str(output_dir / "slide")],
        capture_output=True, text=True
    )

    if result.returncode != 0:
        # Fallback: try using Chrome to screenshot each page
        print("pdftoppm not available. Install with: sudo apt install poppler-utils",
              file=sys.stderr)
        sys.exit(1)

    # Count output files and rename
    pngs = sorted(output_dir.glob("slide-*.png"))
    print(f"Rendered {len(pngs)} pages as PNG images")

    # Create manifest
    import json
    manifest = []
    for i, png in enumerate(pngs, 1):
        manifest.append({
            "slide": i,
            "title": f"Slide {i}",
            "images": [png.name],
        })

    manifest_path = output_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2))
    print(f"Manifest: {manifest_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Extract slide images from NotebookLM output (.pptx or .pdf)"
    )
    parser.add_argument("input", help="Path to .pptx or .pdf file")
    parser.add_argument("-o", "--output", required=True, help="Output directory for extracted images")
    parser.add_argument("--full-slides", action="store_true",
                        help="Render full slides as images (not just embedded images)")

    args = parser.parse_args()
    input_path = Path(args.input)

    if not input_path.exists():
        print(f"Error: {input_path} does not exist", file=sys.stderr)
        sys.exit(1)

    if input_path.suffix.lower() == ".pptx":
        if args.full_slides:
            render_pptx_full_slides(str(input_path), args.output)
        else:
            extract_pptx_images(str(input_path), args.output)
    elif input_path.suffix.lower() == ".pdf":
        render_pdf_pages(str(input_path), args.output)
    else:
        print(f"Unsupported format: {input_path.suffix}. Use .pptx or .pdf", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
